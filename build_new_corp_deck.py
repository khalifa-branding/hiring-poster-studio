import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.xmlchemy import OxmlElement
from PIL import Image, ImageEnhance

# Brand Colors (Aligned with Official Trescon Brandbook)
HEX_DARK_BG = RGBColor(1, 55, 61)         # #01373D (Deep Teal)
HEX_LIGHT_BG = RGBColor(246, 250, 250)    # #F6FAFA (Off-White / Light BG)
HEX_TEAL = RGBColor(0, 165, 163)          # #00A5A3 (Primary Teal)
HEX_DARK_TEAL = RGBColor(1, 55, 61)       # #01373D (Deep Teal for headings on light slides)
HEX_SLATE = RGBColor(70, 77, 83)           # #464D53 (Charcoal / Muted Text)
HEX_ICE_BLUE = RGBColor(230, 239, 240)     # #E6EFF0 (Ice Blue / Soft Accent Text)
HEX_NEON_LIME = RGBColor(192, 244, 60)     # #C0F43C (Electric Lime / Highlights)
HEX_WHITE = RGBColor(255, 255, 255)
HEX_BLACK = RGBColor(0, 0, 0)
HEX_GRAY_BORDER = RGBColor(230, 239, 240) # #E6EFF0 (Soft border)
HEX_BLACK_CARD = RGBColor(30, 33, 36)      # #1E2124 (Charcoal / Dark Card BG)

# Typography (Aligned with Official Trescon Brandbook)
FONT_TITLE = "Anek Devanagari"
FONT_BODY = "Manrope"

# Workspace paths
WORKSPACE_DIR = r"c:\Users\Khalifat\Desktop\Antigravity Projects\Trescon_Corporate_Profile"
IMAGES_DIR = os.path.join(WORKSPACE_DIR, "brand_assets")
SCRATCH_DIR = r"C:\Users\Khalifat\.gemini\antigravity\brain\de77bdf7-0dcd-4f92-8e49-7189290611b4\scratch"
LOGO_PATH = os.path.join(IMAGES_DIR, "10-years-trescon-logo.png")

# Ensure scratch directory exists
os.makedirs(SCRATCH_DIR, exist_ok=True)

# Helper Functions
def create_card(slide, left, top, width, height, bg_color, border_color=None, border_width=1, roundness=0.04):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.adjustments[0] = roundness
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width)
    else:
        shape.line.color.rgb = bg_color
    return shape

def add_textbox(slide, left, top, width, height, text_paragraphs, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0)
    tf.margin_right = Inches(0)
    tf.margin_top = Inches(0)
    tf.margin_bottom = Inches(0)
    
    for i, p_info in enumerate(text_paragraphs):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
            
        p.text = p_info.get("text", "")
        p.alignment = align
        
        font = p.font
        font.name = p_info.get("font_name", FONT_BODY)
        font.size = p_info.get("size", Pt(11))
        font.color.rgb = p_info.get("color", HEX_SLATE)
        font.bold = p_info.get("bold", False)
        font.italic = p_info.get("italic", False)
        
        if "space_after" in p_info:
            p.space_after = p_info["space_after"]
        if "line_spacing" in p_info:
            p.line_spacing = p_info["line_spacing"]
            
    return txBox

def set_shape_transparency(shape, alpha_val):
    fill = shape.fill
    solid_fill = fill._fill
    srgbClr = solid_fill._solidFill.srgbClr
    if srgbClr is not None:
        alpha = OxmlElement('a:alpha')
        alpha.set('val', str(alpha_val))
        srgbClr.append(alpha)

def set_slide_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_slide_background_image(slide, image_filename):
    img_path = os.path.join(IMAGES_DIR, image_filename)
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(0), Inches(0), Inches(13.333), Inches(7.5))

def enhance_image(filename, contrast=1.3, saturation=1.15):
    in_path = os.path.join(IMAGES_DIR, filename)
    out_path = os.path.join(SCRATCH_DIR, f"enhanced_{filename}")
    
    if not os.path.exists(in_path):
        return in_path
        
    try:
        img = Image.open(in_path)
        enh_c = ImageEnhance.Contrast(img)
        img_c = enh_c.enhance(contrast)
        enh_s = ImageEnhance.Color(img_c)
        img_s = enh_s.enhance(saturation)
        img_s.save(out_path)
        return out_path
    except Exception as e:
        print(f"Error enhancing image {filename}: {e}")
        return in_path

def add_signature_header(slide, number_str, title_str, is_dark=False):
    text_color = HEX_WHITE if is_dark else HEX_DARK_TEAL
    line_color = HEX_TEAL
    
    add_textbox(slide, Inches(1.2), Inches(0.5), Inches(10.933), Inches(0.35), [
        {
            "text": f"{number_str}  {title_str.upper()}",
            "font_name": FONT_TITLE,
            "size": Pt(11),
            "color": text_color,
            "bold": True
        }
    ])
    
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(0.9), Inches(10.933), Inches(0.02))
    sep.fill.solid()
    sep.fill.fore_color.rgb = line_color
    sep.line.fill.background()
    
    node = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(0.87), Inches(0.08), Inches(0.08))
    node.fill.solid()
    node.fill.fore_color.rgb = HEX_NEON_LIME
    node.line.fill.background()

def draw_watermark_anniversary(slide, is_dark):
    logo_name = "10-years-trescon-logo-W.png" if is_dark else "10-years-trescon-logo.png"
    watermark_path = os.path.join(IMAGES_DIR, logo_name)
    if os.path.exists(watermark_path):
        slide.shapes.add_picture(watermark_path, Inches(11.233), Inches(6.55), Inches(0.9), Inches(0.55))

# ==================== SLIDE BUILDERS ====================

def build_slide1_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_cover_bg.jpg")
    
    # Asymmetric Typographic Border & Split Visual Container (No raster image)
    create_card(slide, Inches(5.4), Inches(0), Inches(7.933), Inches(7.5), HEX_BLACK_CARD, border_color=HEX_TEAL, border_width=1.5, roundness=0.0)
    
    placeholder_text = [
        {"text": "[ BRAND VISUAL PLACEHOLDER ]", "font_name": FONT_TITLE, "size": Pt(14), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(12)},
        {"text": "Asset: High-Resolution Event Keynote / Digital Art Graphic\nRecommended Aspect Ratio: 16:9 or Full Height Cover Split\nDouble-click to insert your custom editorial or corporate banner photo here.", "font_name": FONT_BODY, "size": Pt(10), "color": HEX_ICE_BLUE, "line_spacing": 1.4, "space_after": Pt(20)},
        {"text": "Design Tip: Choose high contrast, dark-themed photography that aligns with Deep Teal and Electric Lime colors.", "font_name": FONT_BODY, "size": Pt(9), "color": HEX_SLATE, "italic": True}
    ]
    add_textbox(slide, Inches(6.0), Inches(2.2), Inches(6.8), Inches(3.5), placeholder_text)
        
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.38), Inches(0), Inches(0.02), Inches(7.5))
    sep.fill.solid()
    sep.fill.fore_color.rgb = HEX_TEAL
    sep.line.fill.background()
    
    node = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.35), Inches(2.0), Inches(0.08), Inches(0.08))
    node.fill.solid()
    node.fill.fore_color.rgb = HEX_NEON_LIME
    node.line.fill.background()
    
    logo_file = os.path.join(IMAGES_DIR, "10-years-trescon-logo-W.png")
    if os.path.exists(logo_file):
        slide.shapes.add_picture(logo_file, Inches(1.2), Inches(0.7), Inches(1.47), Inches(0.9))
        
    title_paragraphs = [
        {
            "text": "POWERING GLOBAL INNOVATION\nTHROUGH HIGH-IMPACT\nBUSINESS EVENTS",
            "font_name": FONT_TITLE,
            "size": Pt(30),
            "color": HEX_WHITE,
            "bold": True,
            "space_after": Pt(16),
            "line_spacing": 0.95
        },
        {
            "text": "Connecting Businesses with Opportunities across 120+ countries — from concept to measurable outcomes.\n\nGlobal Business Events & Services",
            "font_name": FONT_BODY,
            "size": Pt(11),
            "color": HEX_NEON_LIME,
            "bold": True,
            "space_after": Pt(48)
        },
        {
            "text": "GLOBAL BUSINESS EVENTS & SERVICES COMPANY  |  DUBAI • BANGALORE • GLOBAL",
            "font_name": FONT_TITLE,
            "size": Pt(8.5),
            "color": HEX_ICE_BLUE,
            "bold": True
        }
    ]
    add_textbox(slide, Inches(1.2), Inches(2.0), Inches(4.0), Inches(4.5), title_paragraphs)

def build_slide2_who_we_are(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "01.", "Who We Are", is_dark=False)
    
    manifesto_paragraphs = [
        {
            "text": "Your Strategic Partner for Business Growth",
            "font_name": FONT_TITLE,
            "size": Pt(26),
            "color": HEX_DARK_TEAL,
            "bold": True,
            "space_after": Pt(14)
        },
        {
            "text": "Trescon is a global business events and services company connecting organisations with decision-makers, governments, investors, and industry leaders worldwide.",
            "font_name": FONT_BODY,
            "size": Pt(12.5),
            "color": HEX_SLATE,
            "line_spacing": 1.3
        }
    ]
    add_textbox(slide, Inches(1.2), Inches(1.2), Inches(5.0), Inches(3.4), manifesto_paragraphs)
    
    # Mission Card
    create_card(slide, Inches(6.8), Inches(1.2), Inches(5.3), Inches(1.6), HEX_ICE_BLUE, border_color=HEX_TEAL, border_width=1.0, roundness=0.03)
    p_mission = [
        {
            "text": "Since 2016, we have been enabling technology adoption, fostering strategic partnerships, and delivering measurable business outcomes through high-impact events, managed conferences, and bespoke engagements.",
            "font_name": FONT_BODY,
            "size": Pt(10.5),
            "color": HEX_DARK_TEAL,
            "line_spacing": 1.25
        }
    ]
    add_textbox(slide, Inches(7.05), Inches(1.3), Inches(4.8), Inches(1.4), p_mission)
    
    # Position Card (Highlight)
    create_card(slide, Inches(6.8), Inches(3.0), Inches(5.3), Inches(1.2), HEX_DARK_TEAL, roundness=0.03)
    create_card(slide, Inches(6.8), Inches(3.0), Inches(0.08), Inches(1.2), HEX_NEON_LIME, roundness=0.0)
    p_position = [
        {
            "text": "One Partner. End-to-End Execution. Global Reach.\nStrategy to ROI — managed under one roof",
            "font_name": FONT_TITLE,
            "size": Pt(13),
            "color": HEX_WHITE,
            "bold": True
        }
    ]
    add_textbox(slide, Inches(7.1), Inches(3.2), Inches(4.8), Inches(0.9), p_position)
    
    # Bottom Accent Board & Typographic Placeholder
    create_card(slide, Inches(0), Inches(4.8), Inches(13.333), Inches(2.7), HEX_DARK_TEAL, roundness=0.0)
    create_card(slide, Inches(0), Inches(4.8), Inches(13.333), Inches(0.06), HEX_NEON_LIME, roundness=0.0)
    
    p_strategy = [
        {"text": "GLOBAL STRATEGY ROOM", "font_name": FONT_TITLE, "size": Pt(12), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(2)},
        {"text": "120+ Markets • Real-time Intelligence", "font_name": FONT_BODY, "size": Pt(11), "color": HEX_WHITE, "bold": True, "space_after": Pt(10)},
        {"text": "[ PHOTO PLACEHOLDER: 1920x388 px - Executive Summit Crowds / Onsite Keynote ]", "font_name": FONT_BODY, "size": Pt(8.5), "color": HEX_ICE_BLUE, "italic": True}
    ]
    add_textbox(slide, Inches(1.2), Inches(5.1), Inches(10.9), Inches(2.0), p_strategy)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide3_at_a_glance(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_dark_bg.jpg")
    add_signature_header(slide, "02.", "Trescon at a Glance", is_dark=True)
    
    # Left Hero Segment
    left_paragraphs = [
        {"text": "By the Numbers", "font_name": FONT_TITLE, "size": Pt(28), "color": HEX_WHITE, "bold": True, "space_after": Pt(16)},
        {"text": "By the Numbers — Measurable Global Impact\n\nTrescon at a glance: Eight years of consistent delivery across governments, enterprises and technology ecosystems.",
         "font_name": FONT_BODY, "size": Pt(12.5), "color": HEX_ICE_BLUE, "line_spacing": 1.3}
    ]
    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(4.5), Inches(4.8), left_paragraphs)
    
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(1.4), Inches(0.01), Inches(5.0))
    div.fill.solid()
    div.fill.fore_color.rgb = HEX_TEAL
    div.line.fill.background()
    
    # 8 stats split into 2 columns of 4 rows each
    stats_col1 = [
        ("500+", "Events Delivered", "Across 120+ cities"),
        ("1M+", "Connections Facilitated", "Qualified introductions"),
        ("250K+", "C-Level Attendees", "Decision makers"),
        ("120+", "Countries Represented", "Global network")
    ]
    stats_col2 = [
        ("100+", "Government Partnerships", "Mins. & authorities"),
        ("5,000+", "Exhibitors", "Enterprise & startups"),
        ("3,000+", "Global Speakers", "Industry leaders"),
        ("250+", "Global Team Members", "Events, sales, tech")
    ]
    
    def draw_stats_column(start_x, stats):
        start_y = Inches(1.4)
        row_h = Inches(1.2)
        for idx, (num, label, desc) in enumerate(stats):
            cy = start_y + idx * row_h
            if idx > 0:
                sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, start_x, cy, Inches(3.0), Inches(0.01))
                sep.fill.solid()
                sep.fill.fore_color.rgb = HEX_TEAL
                sep.line.fill.background()
            
            p_stats = [
                {"text": num, "font_name": FONT_TITLE, "size": Pt(20), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(1)},
                {"text": label, "font_name": FONT_BODY, "size": Pt(10), "color": HEX_WHITE, "bold": True, "space_after": Pt(1)},
                {"text": desc, "font_name": FONT_BODY, "size": Pt(8.5), "color": HEX_ICE_BLUE, "italic": True}
            ]
            add_textbox(slide, start_x, cy + Inches(0.08), Inches(3.0), Inches(1.0), p_stats)
            
    draw_stats_column(Inches(6.6), stats_col1)
    draw_stats_column(Inches(10.0), stats_col2)
    
    draw_watermark_anniversary(slide, is_dark=True)

def build_slide4_why_choose_us(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "03.", "Why Choose Trescon", is_dark=False)
    
    left_paragraphs = [
        {
            "text": "Why Global Organisations Choose Trescon",
            "font_name": FONT_TITLE,
            "size": Pt(26),
            "color": HEX_DARK_TEAL,
            "bold": True,
            "space_after": Pt(14),
            "line_spacing": 1.0
        },
        {
            "text": "From concept to execution, we deliver business outcomes — not just events.\n\n10 interlocked capabilities managing the full event lifecycle under one partner — strategy to measurable ROI.",
            "font_name": FONT_BODY,
            "size": Pt(11.5),
            "color": HEX_SLATE,
            "line_spacing": 1.3,
            "space_after": Pt(20)
        },
        {
            "text": "One partner model: Strategy to measurable ROI.",
            "font_name": FONT_BODY,
            "size": Pt(12.5),
            "color": HEX_TEAL,
            "bold": True
        }
    ]
    add_textbox(slide, Inches(1.2), Inches(1.5), Inches(4.5), Inches(4.8), left_paragraphs)
    
    # 10 lifecycle capabilities represented in Bento grid
    capabilities = [
        ("Strategy & Event Concept", "Market mapping & positioning", "01"),
        ("Marketing & Comms", "Full-funnel demand growth", "06"),
        ("Content & Speakers", "Curated agenda & network", "02"),
        ("PR & Media", "Global communications", "07"),
        ("Delegate & Investors", "Qualified C-level sourcing", "03"),
        ("Operations & Logistics", "Venue & onsite execution", "08"),
        ("Sponsorship Sales", "Enterprise partner engine", "04"),
        ("Customer Success", "White-glove experience", "09"),
        ("Branding & Creative", "Identity, stage & digital", "05"),
        ("Project Management", "Single-threaded delivery", "10")
    ]
    
    card_w = Inches(3.2)
    card_h = Inches(0.9)
    gap_x = Inches(0.2)
    gap_y = Inches(0.1)
    start_x = Inches(6.5)
    start_y = Inches(1.3)
    
    for idx, (name, desc, num) in enumerate(capabilities):
        row = idx // 2
        col = idx % 2
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)
        
        create_card(slide, cx, cy, card_w, card_h, HEX_ICE_BLUE, border_color=HEX_GRAY_BORDER, border_width=1.0, roundness=0.08)
        create_card(slide, cx, cy, Inches(0.06), card_h, HEX_TEAL, roundness=0.0)
        
        p_cap = [
            {"text": f"{num}.", "font_name": FONT_BODY, "size": Pt(8.5), "color": HEX_TEAL, "bold": True, "space_after": Pt(1)},
            {"text": name.upper(), "font_name": FONT_TITLE, "size": Pt(9.0), "color": HEX_DARK_TEAL, "bold": True, "space_after": Pt(1)},
            {"text": desc, "font_name": FONT_BODY, "size": Pt(8.0), "color": HEX_SLATE}
        ]
        add_textbox(slide, cx + Inches(0.18), cy + Inches(0.08), card_w - Inches(0.25), card_h - Inches(0.12), p_cap)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide5_divisions(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "04.", "Our Business Divisions", is_dark=False)
    
    divisions = [
        {
            "num": "Flagship",
            "name": "Signature Events",
            "desc": "Large-scale conferences, expos and summits owned by Trescon — World AI Show, HODL, DATE and more. Proven IP with global communities.",
            "highlight": False
        },
        {
            "num": "Gov & Enterprise",
            "name": "Managed Events",
            "desc": "End-to-end event management for governments, associations and enterprises. From concept, marketing to on-ground execution.",
            "highlight": True
        },
        {
            "num": "ROI Focused",
            "name": "Bespoke Events",
            "desc": "Custom executive engagements built around business objectives — pipeline, market entry, product launch and sales acceleration.",
            "highlight": False
        },
        {
            "num": "Enablement",
            "name": "Education & Training",
            "desc": "Training programmes, workshops and bootcamps enabling technology adoption and workforce skilling at scale with measurable outcomes.",
            "highlight": True
        }
    ]
    
    col_w = Inches(2.55)
    gap_x = Inches(0.24)
    start_x = Inches(1.2)
    
    for idx, div in enumerate(divisions):
        cx = start_x + idx * (col_w + gap_x)
        
        if div["highlight"]:
            create_card(slide, cx, Inches(1.5), col_w, Inches(4.8), HEX_DARK_TEAL, roundness=0.03)
            create_card(slide, cx, Inches(1.5), col_w, Inches(0.08), HEX_NEON_LIME, roundness=0.0)
            num_color = HEX_NEON_LIME
            title_color = HEX_WHITE
            desc_color = HEX_ICE_BLUE
            y_offset = Inches(0.1)
        else:
            create_card(slide, cx, Inches(1.6), col_w, Inches(4.6), HEX_WHITE, border_color=HEX_GRAY_BORDER, border_width=1.0, roundness=0.04)
            create_card(slide, cx, Inches(1.6), col_w, Inches(0.08), HEX_TEAL, roundness=0.0)
            num_color = HEX_TEAL
            title_color = HEX_DARK_TEAL
            desc_color = HEX_SLATE
            y_offset = Inches(0.0)
            
        p_div = [
            {"text": div["num"].upper(), "font_name": FONT_BODY, "size": Pt(9.5), "color": num_color, "bold": True, "space_after": Pt(12)},
            {"text": div["name"], "font_name": FONT_TITLE, "size": Pt(13), "color": title_color, "bold": True, "space_after": Pt(12)},
            {"text": div["desc"], "font_name": FONT_BODY, "size": Pt(9.5), "color": desc_color, "line_spacing": 1.25}
        ]
        add_textbox(slide, cx + Inches(0.2), Inches(1.8) + y_offset, col_w - Inches(0.4), Inches(4.0), p_div)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide6_trusted_by(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_dark_bg.jpg")
    add_signature_header(slide, "05.", "Trusted By", is_dark=True)
    
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10.933), Inches(0.3), [
        {"text": "Trusted by Governments, Enterprises & Global Technology Leaders — Ecosystem of 1M+ connections", "font_name": FONT_TITLE, "size": Pt(14), "color": HEX_NEON_LIME, "bold": True}
    ])
    
    logo_groups = [
        {
            "cat": "GOVERNMENTS",
            "logos": ["DUBAI GOV", "MINISTRY OF AI", "KSA MOC", "BENGALURU GOV", "DUBAI FUTURE FOUNDATION"]
        },
        {
            "cat": "TECH PARTNERS",
            "logos": ["AWS", "ORACLE", "MICROSOFT", "IBM", "NEW RELIC", "HITACHI VANTARA"]
        },
        {
            "cat": "ENTERPRISES",
            "logos": ["EMIRATES NBD", "DP WORLD", "CAREEM", "E&", "INFOSYS"]
        },
        {
            "cat": "ASSOCIATIONS",
            "logos": ["DUBAI FINTECH", "AI COUNCIL", "CYBERSECURITY", "CLOUD COLL.", "SUSTAINABILITY", "WEB3 ALLIANCE"]
        },
        {
            "cat": "MEDIA",
            "logos": ["FORBES", "BLOOMBERG", "KHALEEJ TIMES", "GULF NEWS"]
        }
    ]
    
    start_y = Inches(1.5)
    row_h = Inches(0.95)
    
    for idx, lg in enumerate(logo_groups):
        cy = start_y + idx * row_h
        
        # Category Label
        add_textbox(slide, Inches(1.2), cy + Inches(0.12), Inches(2.2), Inches(0.4), [
            {"text": lg["cat"], "font_name": FONT_TITLE, "size": Pt(9.5), "color": HEX_NEON_LIME, "bold": True}
        ])
        
        # Standardized logo card outline boxes
        card_w = Inches(1.4)
        card_h = Inches(0.48)
        gap_x = Inches(0.16)
        
        for l_idx, logo in enumerate(lg["logos"][:6]):
            lx = Inches(3.6) + l_idx * (card_w + gap_x)
            
            create_card(slide, lx, cy + Inches(0.08), card_w, card_h, HEX_BLACK_CARD, border_color=HEX_SLATE, border_width=0.75, roundness=0.08)
            
            add_textbox(slide, lx, cy + Inches(0.18), card_w, card_h - Inches(0.18), [
                {"text": logo.upper(), "font_name": FONT_BODY, "size": Pt(8.5), "color": HEX_WHITE, "bold": True}
            ], align=PP_ALIGN.CENTER)
            
    draw_watermark_anniversary(slide, is_dark=True)

def build_slide7_signature_portfolio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_dark_bg.jpg")
    add_signature_header(slide, "06.", "Signature Event Portfolio", is_dark=True)
    
    # Left Editorial Column
    left_paragraphs = [
        {"text": "Industry-Leading Platforms Driving Global Conversations", "font_name": FONT_TITLE, "size": Pt(24), "color": HEX_WHITE, "bold": True, "space_after": Pt(16)},
        {"text": "Driving conversations across AI, Web3, Cybersecurity, Cloud, Sustainability, Manufacturing and Digital Transformation.",
         "font_name": FONT_BODY, "size": Pt(12.5), "color": HEX_ICE_BLUE, "line_spacing": 1.3}
    ]
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(4.5), Inches(4.5), left_paragraphs)
    
    div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.1), Inches(1.4), Inches(0.01), Inches(5.0))
    div.fill.solid()
    div.fill.fore_color.rgb = HEX_TEAL
    div.line.fill.background()
    
    # Grid of 8 Brands
    brands = [
        ("World AI Show", "AI", "Flagship AI summit, 50+ editions."),
        ("HODL", "Web3", "Longest-running Web3 invest series."),
        ("DATE", "Digital", "Digital acceleration expo for CXOs."),
        ("World Cloud Show", "Cloud", "Cloud adoption & infrastructure."),
        ("World Cyber Security", "Security", "CISO-focused risk & resilience."),
        ("Future Sustainability", "ESG", "Net-zero & ESG with govts."),
        ("Future Factory Show", "Manufacturing", "Smart mfg & Industry 4.0."),
        ("CARE", "HealthTech", "Healthcare innovation & care models.")
    ]
    
    card_w = Inches(2.9)
    card_h = Inches(1.05)
    gap_x = Inches(0.2)
    gap_y = Inches(0.2)
    start_x = Inches(6.5)
    start_y = Inches(1.6)
    
    for idx, (title, cat, desc) in enumerate(brands):
        row = idx // 2
        col = idx % 2
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)
        
        create_card(slide, cx, cy, card_w, card_h, HEX_BLACK_CARD, border_color=HEX_TEAL, border_width=1.0, roundness=0.08)
        
        p_brand = [
            {"text": cat.upper(), "font_name": FONT_BODY, "size": Pt(8.5), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(2)},
            {"text": title, "font_name": FONT_TITLE, "size": Pt(11), "color": HEX_WHITE, "bold": True, "space_after": Pt(2)},
            {"text": desc, "font_name": FONT_BODY, "size": Pt(8.0), "color": HEX_ICE_BLUE}
        ]
        add_textbox(slide, cx + Inches(0.18), cy + Inches(0.12), card_w - Inches(0.3), card_h - Inches(0.2), p_brand)
        
    draw_watermark_anniversary(slide, is_dark=True)

def build_slide8_managed_portfolio(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "07.", "Managed Events Portfolio", is_dark=False)
    
    left_paragraphs = [
        {"text": "Trusted by Governments and Leading Organisations to Deliver Flagships", "font_name": FONT_TITLE, "size": Pt(22), "color": HEX_DARK_TEAL, "bold": True, "space_after": Pt(16)},
        {"text": "Trusted to conceptualise, market and execute flagship events with C-level acquisition, government protocol and P&L ownership.",
         "font_name": FONT_BODY, "size": Pt(12.5), "color": HEX_SLATE, "line_spacing": 1.3}
    ]
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(4.2), Inches(4.5), left_paragraphs)
    
    # 6 Managed Events
    events = [
        ("Dubai FinTech Summit", "Dubai • DIFC", "World's largest FinTech gathering under Dubai government | Finance • Government"),
        ("Future Sustainability Forum", "UAE", "Sustainability & ESG leadership with ministries & funds | ESG • Net Zero"),
        ("Dubai AI Festival", "Dubai AI Campus", "AI adoption festival powering enterprise transformation | AI • Festival"),
        ("World Police Summit", "Dubai Police", "Global law enforcement technology and security innovation summit | Security • Gov"),
        ("Bengaluru Skill Summit", "Karnataka Gov", "Skilling & future of work summit enabling youth employability | Skilling • India"),
        ("Dubai Future Forum", "Dubai Future Foundation", "Futures thinking and moonshot innovation with DFF | Future • Innovation")
    ]
    
    card_w = Inches(3.3)
    card_h = Inches(1.5)
    gap_x = Inches(0.24)
    gap_y = Inches(0.2)
    start_x = Inches(5.8)
    start_y = Inches(1.6)
    
    for idx, (title, subtitle, desc) in enumerate(events):
        row = idx // 2
        col = idx % 2
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)
        
        create_card(slide, cx, cy, card_w, card_h, HEX_ICE_BLUE, border_color=HEX_GRAY_BORDER, border_width=1.0, roundness=0.04)
        create_card(slide, cx, cy, Inches(0.08), card_h, HEX_TEAL, roundness=0.0)
        
        p_ev = [
            {"text": f"{title}  |  {subtitle}", "font_name": FONT_TITLE, "size": Pt(11), "color": HEX_DARK_TEAL, "bold": True, "space_after": Pt(6)},
            {"text": desc, "font_name": FONT_BODY, "size": Pt(9.5), "color": HEX_SLATE, "line_spacing": 1.25}
        ]
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.12), card_w - Inches(0.35), card_h - Inches(0.2), p_ev)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide9_bespoke_events(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "08.", "Bespoke Events", is_dark=False)
    
    add_textbox(slide, Inches(1.2), Inches(1.2), Inches(10.933), Inches(0.5), [
        {"text": "Bespoke • ROI Led — Custom Events Designed Around Your Business Goals", "font_name": FONT_TITLE, "size": Pt(22), "color": HEX_DARK_TEAL, "bold": True}
    ])
    
    # Left Content Column
    p_left = [
        {"text": "We help organisations:", "font_name": FONT_TITLE, "size": Pt(13), "color": HEX_DARK_TEAL, "bold": True, "space_after": Pt(12)},
        {"text": "• Launch products to target C-level buyers\n• Generate qualified pipelines and accelerate sales\n• Build thought leadership in emerging sectors\n• Enter new global technology corridors\n• Engage key state & enterprise decision-makers",
         "font_name": FONT_BODY, "size": Pt(11), "color": HEX_SLATE, "line_spacing": 1.5, "space_after": Pt(24)},
        {"text": "Every engagement is built around measurable business outcomes — not vanity metrics.", "font_name": FONT_BODY, "size": Pt(11), "color": HEX_TEAL, "bold": True}
    ]
    add_textbox(slide, Inches(1.2), Inches(1.9), Inches(5.0), Inches(4.5), p_left)
    
    # Right Image Card
    img_path = enhance_image("page_26_img_1_1920x1080.png")
    
    # Stats overlay
    create_card(slide, Inches(6.8), Inches(1.9), Inches(5.333), Inches(1.0), HEX_DARK_TEAL, roundness=0.03)
    p_stats = [
        {
            "text": "27-70 meetings/ed.  |  ~68% C-level  |  120+ countries",
            "font_name": FONT_TITLE,
            "size": Pt(12),
            "color": HEX_NEON_LIME,
            "bold": True,
            "space_after": Pt(2)
        },
        {
            "text": "Outcome first: Pipeline, pilots and partnership KPIs",
            "font_name": FONT_BODY,
            "size": Pt(10),
            "color": HEX_WHITE
        }
    ]
    add_textbox(slide, Inches(7.0), Inches(2.0), Inches(5.0), Inches(0.8), p_stats)
    
    # Right-side Typographic Visual Placeholder (No raster image)
    create_card(slide, Inches(6.8), Inches(3.0), Inches(5.333), Inches(3.1), HEX_LIGHT_BG, border_color=HEX_TEAL, border_width=1.5, roundness=0.04)
    
    placeholder_text_s9 = [
        {"text": "[ B2B KEYNOTE IMAGE PLACEHOLDER ]", "font_name": FONT_TITLE, "size": Pt(11), "color": HEX_DARK_TEAL, "bold": True, "space_after": Pt(6)},
        {"text": "Asset: Closed-door VIP Boardroom or Matchmaking Session Photo\nSuggested Size: 5.3 x 3.1 Inches\nDouble-click in PowerPoint to swap with your actual bespoke event networking visual.", "font_name": FONT_BODY, "size": Pt(9.5), "color": HEX_SLATE, "line_spacing": 1.3}
    ]
    add_textbox(slide, Inches(7.1), Inches(3.4), Inches(4.8), Inches(2.3), placeholder_text_s9)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide10_why_trescon(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "09.", "Why Trescon", is_dark=False)
    
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10.933), Inches(0.4), [
        {"text": "Why Clients Work With Trescon", "font_name": FONT_TITLE, "size": Pt(20), "color": HEX_DARK_TEAL, "bold": True}
    ])
    
    benefits = [
        ("Reach • Relevance", "Global network across 120+ countries. Built through 500+ events and government alliances. Access where your buying committee sits."),
        ("68% C-level avg", "Access to qualified C-level decision-makers. 250k+ C-level attendees historically, with qualification and intent mapping."),
        ("Trust • Access", "Government and enterprise relationships. 100+ government partnerships enabling protocol, policy context and trusted convening power."),
        ("Zero compromise", "Proven event delivery capability. Stage production, content curation, logistics and hospitality — 8 years of flagship-grade execution."),
        ("One P&L", "End-to-end execution under one partner. Strategy, sales, marketing, PR, ops and success — single-threaded ownership from concept to ROI."),
        ("Measurable Outcomes", "Results-driven approach backed by measurable KPIs. Pipeline meetings, qualified attendees, pilots booked — defined upfront and reported weekly.")
    ]
    
    card_w = Inches(3.45)
    card_h = Inches(2.0)
    gap_x = Inches(0.29)
    gap_y = Inches(0.24)
    start_x = Inches(1.2)
    start_y = Inches(1.6)
    
    for idx, (title, desc) in enumerate(benefits):
        row = idx // 3
        col = idx % 3
        cx = start_x + col * (card_w + gap_x)
        cy = start_y + row * (card_h + gap_y)
        
        # Highlight the second block
        if idx == 1:
            create_card(slide, cx, cy, card_w, card_h, HEX_DARK_TEAL, roundness=0.03)
            create_card(slide, cx, cy, Inches(0.08), card_h, HEX_NEON_LIME, roundness=0.0)
            t_color = HEX_NEON_LIME
            desc_color = HEX_ICE_BLUE
        else:
            create_card(slide, cx, cy, card_w, card_h, HEX_ICE_BLUE, border_color=HEX_GRAY_BORDER, border_width=1.0, roundness=0.04)
            create_card(slide, cx, cy, Inches(0.08), card_h, HEX_TEAL, roundness=0.0)
            t_color = HEX_DARK_TEAL
            desc_color = HEX_SLATE
            
        p_ben = [
            {"text": f"0{idx+1}.", "font_name": FONT_TITLE, "size": Pt(12), "color": HEX_NEON_LIME if idx==1 else HEX_TEAL, "bold": True, "space_after": Pt(4)},
            {"text": title, "font_name": FONT_TITLE, "size": Pt(11.5), "color": HEX_WHITE if idx==1 else HEX_DARK_TEAL, "bold": True, "space_after": Pt(6)},
            {"text": desc, "font_name": FONT_BODY, "size": Pt(9.0), "color": desc_color, "line_spacing": 1.2}
        ]
        add_textbox(slide, cx + Inches(0.2), cy + Inches(0.15), card_w - Inches(0.35), card_h - Inches(0.3), p_ben)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide11_success_stories(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "10.", "Selected Success Stories", is_dark=False)
    
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10.933), Inches(0.4), [
        {"text": "Delivering Measurable Business Outcomes", "font_name": FONT_TITLE, "size": Pt(22), "color": HEX_DARK_TEAL, "bold": True}
    ])
    
    cases = [
        {
            "brand": "AWS (ASEAN)",
            "obj": "Expand Greenfield cloud adoption across ASEAN.",
            "res": "27 qualified enterprise meetings, Pilot programme opportunities, Increased enterprise engagement.\n\nMetrics: 27 meetings • ASEAN market",
            "highlight": False
        },
        {
            "brand": "New Relic (India)",
            "obj": "Build market awareness in India.",
            "res": "70 qualified attendees, Product-market validation, Strong GTM business pipeline.\n\nMetrics: 70 qualified • India GTM",
            "highlight": True
        },
        {
            "brand": "Hitachi Vantara + Oracle",
            "obj": "Grow strategic partnerships.",
            "res": "39 qualified enterprise attendees, Stronger partner ecosystem, Increased solution awareness.\n\nMetrics: 39 enterprise • 2x partners",
            "highlight": False
        }
    ]
    
    col_w = Inches(3.45)
    gap_x = Inches(0.29)
    start_x = Inches(1.2)
    
    for idx, c in enumerate(cases):
        cx = start_x + idx * (col_w + gap_x)
        
        if c["highlight"]:
            create_card(slide, cx, Inches(1.5), col_w, Inches(4.8), HEX_DARK_TEAL, roundness=0.03)
            create_card(slide, cx, Inches(1.5), col_w, Inches(0.08), HEX_NEON_LIME, roundness=0.0)
            t_color = HEX_NEON_LIME
            body_color = HEX_WHITE
            desc_color = HEX_ICE_BLUE
            y_offset = Inches(0.1)
        else:
            create_card(slide, cx, Inches(1.6), col_w, Inches(4.6), HEX_WHITE, border_color=HEX_GRAY_BORDER, border_width=1.0, roundness=0.04)
            create_card(slide, cx, Inches(1.6), col_w, Inches(0.08), HEX_TEAL, roundness=0.0)
            t_color = HEX_TEAL
            body_color = HEX_DARK_TEAL
            desc_color = HEX_SLATE
            y_offset = Inches(0.0)
            
        p_case = [
            {"text": c["brand"].upper(), "font_name": FONT_TITLE, "size": Pt(14), "color": HEX_WHITE if c["highlight"] else HEX_DARK_TEAL, "bold": True, "space_after": Pt(12)},
            
            {"text": "OBJECTIVE", "font_name": FONT_TITLE, "size": Pt(8.5), "color": t_color, "bold": True, "space_after": Pt(2)},
            {"text": c["obj"], "font_name": FONT_BODY, "size": Pt(10), "color": desc_color, "line_spacing": 1.25, "space_after": Pt(16)},
            
            {"text": "RESULT & IMPACT", "font_name": FONT_TITLE, "size": Pt(8.5), "color": t_color, "bold": True, "space_after": Pt(2)},
            {"text": c["res"], "font_name": FONT_BODY, "size": Pt(9.5), "color": desc_color, "line_spacing": 1.25, "bold": True}
        ]
        add_textbox(slide, cx + Inches(0.2), Inches(1.8) + y_offset, col_w - Inches(0.4), Inches(4.0), p_case)
        
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide12_upcoming_events(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_dark_bg.jpg")
    add_signature_header(slide, "11.", "Upcoming Flagship Events", is_dark=True)
    
    add_textbox(slide, Inches(1.2), Inches(1.0), Inches(10.933), Inches(0.3), [
        {"text": "Meet us across the world's leading innovation ecosystems. Calendar is dynamic — confirm on site.", "font_name": FONT_BODY, "size": Pt(11.5), "color": HEX_NEON_LIME, "bold": True}
    ])
    
    upcoming = [
        ("Dubai FinTech Summit", "May 2025 • Dubai\nFinTech | Flagship", "Madinat Jumeirah, Dubai, UAE"),
        ("Future Sustainability", "June 2025 • Dubai\nESG | Gov-led", "Museum of Future, Dubai, UAE"),
        ("Future Islamic Finance", "Q3 2025 • GCC\nFinance | New", "Dubai / Riyadh, GCC"),
        ("Bengaluru Skill Summit", "Q3 2025 • India\nSkilling | India", "BIEC, Bengaluru, India"),
        ("Future Factory Show", "Q4 2025 • Global\nIndustrial | Global", "Dubai • Berlin, Global")
    ]
    
    card_w = Inches(2.0)
    gap_x = Inches(0.233)
    start_x = Inches(1.2)
    
    for idx, (title, sub, location) in enumerate(upcoming):
        cx = start_x + idx * (card_w + gap_x)
        
        # Sleek scheduler card
        create_card(slide, cx, Inches(1.8), card_w, Inches(4.5), HEX_BLACK_CARD, border_color=HEX_TEAL, border_width=1.0, roundness=0.08)
        create_card(slide, cx, Inches(1.8), card_w, Inches(0.06), HEX_NEON_LIME, roundness=0.0)
        
        p_up = [
            {"text": f"0{idx+1}", "font_name": FONT_TITLE, "size": Pt(18), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(16)},
            {"text": title.upper(), "font_name": FONT_TITLE, "size": Pt(11), "color": HEX_WHITE, "bold": True, "space_after": Pt(12), "line_spacing": 1.1},
            {"text": "CATEGORY", "font_name": FONT_BODY, "size": Pt(8), "color": HEX_TEAL, "bold": True, "space_after": Pt(2)},
            {"text": sub, "font_name": FONT_BODY, "size": Pt(9), "color": HEX_ICE_BLUE, "space_after": Pt(16)},
            {"text": "LOCATION", "font_name": FONT_BODY, "size": Pt(8), "color": HEX_TEAL, "bold": True, "space_after": Pt(2)},
            {"text": location, "font_name": FONT_BODY, "size": Pt(9), "color": HEX_ICE_BLUE, "bold": True}
        ]
        add_textbox(slide, cx + Inches(0.15), Inches(2.1), card_w - Inches(0.3), Inches(4.0), p_up)
        
    draw_watermark_anniversary(slide, is_dark=True)

def build_slide13_leadership(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_light_bg.jpg")
    add_signature_header(slide, "12.", "Leadership", is_dark=False)
    
    leaders = [
        ("FOUNDER &\nCHAIRMAN", "Vision, Gov Relations & Global Growth | Leadership"),
        ("CO-FOUNDER &\nVICE CHAIRMAN", "Strategy & Business Development | Growth"),
        ("GROUP CEO", "Operations, Product & Commercial | Executive"),
        ("MANAGING\nDIRECTOR", "Delivery, Marketing & Partnerships | Delivery")
    ]
    
    col_w = Inches(2.55)
    gap_x = Inches(0.24)
    start_x = Inches(1.2)
    
    for idx, (role, desc) in enumerate(leaders):
        cx = start_x + idx * (col_w + gap_x)
        
        # Profile outline card
        create_card(slide, cx, Inches(1.5), col_w, Inches(3.2), HEX_WHITE, border_color=HEX_GRAY_BORDER, border_width=1.0, roundness=0.04)
        create_card(slide, cx, Inches(1.5), col_w, Inches(0.06), HEX_TEAL, roundness=0.0)
        
        p_leader = [
            {"text": f"LEADER 0{idx+1}", "font_name": FONT_BODY, "size": Pt(8.5), "color": HEX_TEAL, "bold": True, "space_after": Pt(8)},
            {"text": role, "font_name": FONT_TITLE, "size": Pt(13), "color": HEX_DARK_TEAL, "bold": True, "space_after": Pt(10), "line_spacing": 1.05},
            {"text": desc, "font_name": FONT_BODY, "size": Pt(10.0), "color": HEX_SLATE, "line_spacing": 1.25}
        ]
        add_textbox(slide, cx + Inches(0.18), Inches(1.7), col_w - Inches(0.35), Inches(2.8), p_leader)
        
    # Bottom Anchor Card
    create_card(slide, Inches(1.2), Inches(4.9), Inches(10.933), Inches(1.3), HEX_DARK_TEAL, roundness=0.03)
    create_card(slide, Inches(1.2), Inches(4.9), Inches(0.08), Inches(1.3), HEX_NEON_LIME, roundness=0.0)
    
    p_anchor = [
        {"text": "GLOBAL TALENT INFRASTRUCTURE  |  BUILT FOR C-LEVEL OUTCOMES", "font_name": FONT_TITLE, "size": Pt(9.5), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(4)},
        {"text": "Supported by 250+ professionals worldwide across sales, content, design, ops, marketing and customer success. Every function runs as a revenue center with pipeline targets.",
         "font_name": FONT_BODY, "size": Pt(11), "color": HEX_WHITE, "line_spacing": 1.2, "bold": True, "space_after": Pt(4)},
        {"text": "Specialties: Event Ops • Marketing • Commercial • Production • Tech • Content", "font_name": FONT_BODY, "size": Pt(9.5), "color": HEX_ICE_BLUE}
    ]
    add_textbox(slide, Inches(1.6), Inches(5.05), Inches(10.0), Inches(1.05), p_anchor)
    
    draw_watermark_anniversary(slide, is_dark=False)

def build_slide14_closing(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_slide_background_image(slide, "trescon_closing_bg.jpg")
    
    # Clean Typographic Solid Color Background (No raster images)
    
    # 10-Year Anniversary Logo placement
    logo_file = os.path.join(IMAGES_DIR, "10-years-trescon-logo-W.png")
    if os.path.exists(logo_file):
        slide.shapes.add_picture(logo_file, Inches(1.2), Inches(5.6), Inches(1.47), Inches(0.9))
        
    closing_paragraphs = [
        {
            "text": "Ready to Scale with Trescon — Let's Build Your Next High-Impact Event",
            "font_name": FONT_TITLE,
            "size": Pt(32),
            "color": HEX_WHITE,
            "bold": True,
            "space_after": Pt(16),
            "line_spacing": 0.98
        },
        {
            "text": "Whether you're launching a flagship summit, expanding into new markets, or engaging enterprise buyers, Trescon delivers the strategy, audience and execution to make it happen.",
            "font_name": FONT_BODY,
            "size": Pt(12.5),
            "color": HEX_ICE_BLUE,
            "line_spacing": 1.3,
            "space_after": Pt(22)
        },
        {
            "text": "CONNECTING BUSINESSES WITH OPPORTUNITIES",
            "font_name": FONT_TITLE,
            "size": Pt(10),
            "color": HEX_NEON_LIME,
            "bold": True
        }
    ]
    add_textbox(slide, Inches(1.2), Inches(1.8), Inches(8.5), Inches(3.5), closing_paragraphs)
    
    contact_paragraphs = [
        {"text": "GLOBAL HEADQUARTERS", "font_name": FONT_TITLE, "size": Pt(8.5), "color": HEX_NEON_LIME, "bold": True, "space_after": Pt(2)},
        {"text": "Liberty House, DIFC, Dubai, UAE", "font_name": FONT_BODY, "size": Pt(10), "color": HEX_WHITE, "space_after": Pt(2)},
        {"text": "info@tresconglobal.com  |  www.tresconglobal.com", "font_name": FONT_BODY, "size": Pt(10), "color": HEX_ICE_BLUE}
    ]
    add_textbox(slide, Inches(8.2), Inches(5.6), Inches(3.9), Inches(1.2), contact_paragraphs, align=PP_ALIGN.RIGHT)

def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    build_slide1_cover(prs)
    build_slide2_who_we_are(prs)
    build_slide3_at_a_glance(prs)
    build_slide4_why_choose_us(prs)
    build_slide5_divisions(prs)
    build_slide6_trusted_by(prs)
    build_slide7_signature_portfolio(prs)
    build_slide8_managed_portfolio(prs)
    build_slide9_bespoke_events(prs)
    build_slide10_why_trescon(prs)
    build_slide11_success_stories(prs)
    build_slide12_upcoming_events(prs)
    build_slide13_leadership(prs)
    build_slide14_closing(prs)
    
    out_file = os.path.join(WORKSPACE_DIR, "Trescon_Corporate_Pitch_Deck_v7.pptx")
    prs.save(out_file)
    print(f"Presentation generated and saved to: {out_file}")

if __name__ == "__main__":
    build_presentation()
